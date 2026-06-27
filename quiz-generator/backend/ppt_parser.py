"""
PPT Parser Module
Extracts text content from .pptx files
"""

from pptx import Presentation
import os

def extract_text_from_pptx(file_path):
    """
    Extract all text content from a .pptx file.
    
    Args:
        file_path (str): Path to the .pptx file
        
    Returns:
        dict: Contains slide_count and slides (list of slide texts)
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    if not file_path.lower().endswith('.pptx'):
        raise ValueError("File must be a .pptx file")
    
    try:
        prs = Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(row_text))
            
            slides_content.append({
                "slide_number": i,
                "text": "\n".join(slide_text) if slide_text else "(No text content)"
            })
        
        return {
            "slide_count": len(slides_content),
            "slides": slides_content
        }
    
    except Exception as e:
        raise Exception(f"Error parsing PPTX file: {str(e)}")


def get_combined_text(file_path):
    """
    Extract and combine all text from a PPTX file into a single string.
    
    Args:
        file_path (str): Path to the .pptx file
        
    Returns:
        str: Combined text from all slides
    """
    result = extract_text_from_pptx(file_path)
    combined = []
    for slide in result["slides"]:
        combined.append(f"--- Slide {slide['slide_number']} ---\n{slide['text']}")
    
    return "\n\n".join(combined)